"""
engine.py - Moteur de rendu deterministe (schema JSON -> .pptx natif charte OZ).

Coeur reutilisable, sans IA. Multi-theme (clair / sombre). Appele par le serveur
MCP (Cowork) comme par un serveur Claude SDK. La charte vit dans theme.py.
"""

from __future__ import annotations
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData

from .theme import Brand, Theme, get_theme
from . import icons


# ----------------------------- helpers -----------------------------

def _no_line(shape):
    shape.line.fill.background()


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    _no_line(shape)


def _round(shape, radius=0.06):
    try:
        shape.adjustments[0] = radius
    except Exception:
        pass


def _text(slide, txt, left, top, width, height, size, color,
          bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          line_spacing=1.3, caps=False, italic=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = txt.upper() if caps else txt
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.name = Brand.FONT; f.color.rgb = color
    return box


def _first_word_bold(p, text, size, color):
    parts = text.split(" ", 1)
    r1 = p.add_run(); r1.text = parts[0]
    r1.font.bold = True; r1.font.size = Pt(size); r1.font.name = Brand.FONT; r1.font.color.rgb = color
    if len(parts) > 1:
        r2 = p.add_run(); r2.text = " " + parts[1]
        r2.font.bold = False; r2.font.size = Pt(size); r2.font.name = Brand.FONT; r2.font.color.rgb = color


def _gradient_bg(slide, c1, c2, angle=45):
    fill = slide.background.fill
    fill.gradient()
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[1].position = 1.0
    fill.gradient_stops[1].color.rgb = c2
    try:
        fill.gradient_angle = angle
    except Exception:
        pass


def _solid_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


# --------------------------- composants ---------------------------

_ASSETS_DIR = os.path.join(os.path.dirname(__file__), "assets")
# Ratios largeur/hauteur des logos officiels
_LOGO_RATIO = {"principal": 1.0, "secondaire": 1500 / 1364}


def _logo_path(kind: str, on_dark: bool) -> str:
    variant = "blanc" if on_dark else "orange"   # charte : blanc/sombre, orange/clair
    return os.path.join(_ASSETS_DIR, f"logo_{kind}_{variant}.png")


def _logo(slide, left, top, height, th: Theme, kind: str = "secondaire"):
    """Pose le logo officiel (principal = avec wordmark ; secondaire = badge seul)."""
    path = _logo_path(kind, th.on_dark)
    width = height * _LOGO_RATIO.get(kind, 1.0)
    slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
    return width


def _footer(slide, th: Theme, page_no=None, total=None):
    line_col = th.text_muted if th.on_dark else Brand.FILET
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(Brand.MARGIN), Inches(Brand.FOOTER_TOP),
                                Inches(Brand.SW - 2*Brand.MARGIN), Pt(0.75))
    _fill(ln, line_col)
    _logo(slide, Brand.MARGIN, Brand.FOOTER_TOP + 0.11, 0.26, th, "secondaire")
    _text(slide, Brand.FOOTER_LEFT, Brand.MARGIN + 0.44, Brand.FOOTER_TOP + 0.12, 6, 0.3, 9,
          th.text_muted, caps=True)
    right = Brand.FOOTER_RIGHT
    if page_no is not None and total is not None:
        right = f"{Brand.FOOTER_RIGHT}     PAGE {page_no} / {total}"
    _text(slide, right, Brand.SW - Brand.MARGIN - 6, Brand.FOOTER_TOP + 0.12, 6, 0.3, 9,
          th.text_muted, align=PP_ALIGN.RIGHT, caps=True)


def _title(slide, title, th: Theme, subtitle=None):
    # Charte : H1 Helvetica Bold 24pt capitales. Barre verticale orange (5px)
    # alignee sur la hauteur des capitales du titre, centree verticalement.
    box_h = 0.5
    center = Brand.TITLE_TOP + box_h / 2
    bar_h = 0.34  # ~ hauteur des capitales a 24pt
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(Brand.MARGIN),
                                 Inches(center - bar_h / 2), Pt(5), Inches(bar_h))
    _fill(bar, Brand.ORANGE)
    _text(slide, title, Brand.MARGIN + 0.2, Brand.TITLE_TOP, Brand.SW - 2*Brand.MARGIN, box_h,
          24, th.text, bold=True, caps=True, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        # H2 : Helvetica Bold 14pt, orange, sous le titre
        _text(slide, subtitle, Brand.MARGIN + 0.2, Brand.TITLE_TOP + box_h, Brand.SW - 2*Brand.MARGIN,
              0.35, 14, Brand.ORANGE, bold=True)


def _icon_square(slide, left, top, size, name, th: Theme, orange=False):
    sq = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                Inches(size), Inches(size))
    _fill(sq, th.tile_icon_bg); _round(sq, 0.15)
    col = icons.hex_to_rgba("EC4324") if orange else icons.hex_to_rgba(th.icon_hex)
    icon = size * 0.6; off = (size - icon) / 2
    slide.shapes.add_picture(icons.icon_png(name, col), Inches(left + off), Inches(top + off),
                             Inches(icon), Inches(icon))


def _tile(slide, left, top, width, height, th: Theme, heading=None, body=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    _fill(card, th.tile); _round(card)
    pad = 0.28; y = top + pad
    if heading:
        _text(slide, heading, left + pad, y, width - 2*pad, 0.4, 14, Brand.ORANGE, bold=True)
        y += 0.5
    if body:
        box = slide.shapes.add_textbox(Inches(left + pad), Inches(y),
                                       Inches(width - 2*pad), Inches(height - (y - top) - pad))
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]; p.line_spacing = 1.4; p.alignment = PP_ALIGN.LEFT
        _first_word_bold(p, body, 12.5, th.text)


# ----------------------------- layouts -----------------------------

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def layout_cover(prs, c, th):
    s = _blank(prs); _gradient_bg(s, th.bg_grad_a, th.bg_grad_b)
    if c.get("kicker"):
        _text(s, c["kicker"], Brand.MARGIN, 2.4, 9, 0.4, 13, th.text_muted, caps=True)
    _text(s, c["title"], Brand.MARGIN, 2.7, 10.5, 1.6, 46, th.text, bold=True, caps=True, line_spacing=1.05)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(Brand.MARGIN), Inches(4.35), Inches(0.5), Pt(3))
    _fill(bar, Brand.ORANGE)
    if c.get("subtitle"):
        _text(s, c["subtitle"], Brand.MARGIN, 4.55, 10, 0.6, 15, th.text if th.on_dark else th.text_muted)
    return s


def layout_section(prs, c, th):
    s = _blank(prs); _gradient_bg(s, th.bg_grad_a, th.bg_grad_b)
    if c.get("number"):
        _text(s, c["number"], Brand.MARGIN, 2.6, 3, 1.2, 66, Brand.ORANGE, bold=True)
    _text(s, c["title"], Brand.MARGIN, 3.7, 11, 1.2, 34, th.text, bold=True, caps=True)
    return s


def layout_content(prs, c, th):
    s = _blank(prs); _solid_bg(s, th.bg); _title(s, c["title"], th)
    y = 1.9
    for b in c.get("blocks", []):
        _text(s, b["heading"], Brand.MARGIN, y, 11, 0.4, 14, Brand.ORANGE, bold=True)
        box = s.shapes.add_textbox(Inches(Brand.MARGIN), Inches(y + 0.42),
                                   Inches(Brand.SW - 2*Brand.MARGIN), Inches(0.9))
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]; p.line_spacing = 1.4
        _first_word_bold(p, b["body"], 13, th.text)
        y += 1.25
    return s


def layout_tiles(prs, c, th):
    s = _blank(prs); _solid_bg(s, th.bg); _title(s, c["title"], th)
    tiles = c.get("tiles", []); n = max(1, len(tiles)); gap = 0.35
    total_w = Brand.SW - 2*Brand.MARGIN
    tw = (total_w - gap*(n-1)) / n
    for i, t in enumerate(tiles):
        _tile(s, Brand.MARGIN + i*(tw+gap), 2.4, tw, 3.2, th, t.get("heading"), t.get("body"))
    return s


def layout_grid(prs, c, th):
    """Grille de 4 a 6 tuiles courtes sur 2 ou 3 colonnes (2 rangees)."""
    s = _blank(prs); _solid_bg(s, th.bg); _title(s, c["title"], th)
    items = c.get("items", []); n = max(1, len(items))
    top0 = 1.85
    if c.get("intro"):
        _text(s, c["intro"], Brand.MARGIN, 1.55, Brand.SW - 2*Brand.MARGIN, 0.4, 12.5,
              th.text if th.on_dark else th.text_muted)
        top0 = 2.15
    cols = 3 if n > 4 else 2
    import math
    rows = math.ceil(n / cols)
    gap = 0.3
    total_w = Brand.SW - 2*Brand.MARGIN
    tw = (total_w - gap*(cols-1)) / cols
    bottom = Brand.FOOTER_TOP - 0.2
    th_h = (bottom - top0 - gap*(rows-1)) / rows
    for i, it in enumerate(items):
        r, col = divmod(i, cols)
        left = Brand.MARGIN + col*(tw+gap)
        top = top0 + r*(th_h+gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(tw), Inches(th_h))
        _fill(card, th.tile); _round(card, 0.08)
        pad = 0.22; y = top + pad
        if it.get("heading"):
            _text(s, it["heading"], left + pad, y, tw - 2*pad, 0.35, 13, Brand.ORANGE, bold=True, line_spacing=1.05)
            y += 0.42
        if it.get("body"):
            box = s.shapes.add_textbox(Inches(left + pad), Inches(y), Inches(tw - 2*pad), Inches(top + th_h - y - pad*0.5))
            tf = box.text_frame; tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]; p.line_spacing = 1.25
            _first_word_bold(p, it["body"], 11, th.text)
    return s


def layout_kpi(prs, c, th):
    s = _blank(prs); _solid_bg(s, th.bg); _title(s, c["title"], th)
    kpis = c.get("kpis", []); n = max(1, len(kpis)); gap = 0.3
    total_w = Brand.SW - 2*Brand.MARGIN
    cw = (total_w - gap*(n-1)) / n; ch = 2.9; top = 2.6
    for i, k in enumerate(kpis):
        left = Brand.MARGIN + i*(cw+gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(cw), Inches(ch))
        _fill(card, th.tile); _round(card)
        pad = 0.28
        _icon_square(s, left + pad, top + pad, 0.6, k.get("icon", "bar-chart-3"), th, orange=k.get("accent", False))
        _text(s, k.get("label", ""), left + pad, top + 1.05, cw - 2*pad, 0.5, 11, th.text_muted, caps=True, line_spacing=1.15)
        _text(s, k.get("value", ""), left + pad, top + 1.5, cw - 2*pad, 0.7, 32,
              Brand.ORANGE if k.get("accent") else th.text, bold=True)
        if k.get("delta"):
            pos = k.get("positive", True)
            _text(s, k["delta"], left + pad, top + 2.25, cw - 2*pad, 0.4, 13,
                  Brand.VERT if pos else Brand.ORANGE, bold=True)
    return s


def layout_table(prs, c, th):
    s = _blank(prs); _solid_bg(s, th.bg); _title(s, c["title"], th)
    headers = c["headers"]; rows = c["rows"]
    left_align = c.get("left_align", False)  # tableaux a texte (pas de chiffres)
    nrows = len(rows) + 1; ncols = len(headers)
    height = min(4.6, 0.5*nrows)
    tbl = s.shapes.add_table(nrows, ncols, Inches(Brand.MARGIN), Inches(2.1),
                             Inches(Brand.SW - 2*Brand.MARGIN), Inches(height)).table
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = Brand.NOIR
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if left_align else (PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT)
        r = p.add_run(); r.text = str(h).upper()
        r.font.bold = True; r.font.size = Pt(10); r.font.name = Brand.FONT; r.font.color.rgb = Brand.BLANC
    for i, row in enumerate(rows):
        band = th.table_body_alt if i % 2 else th.table_body_base
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = band
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if left_align else (PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT)
            r = p.add_run(); txt = str(val)
            r.font.size = Pt(10); r.font.name = Brand.FONT
            if txt.startswith("!"):
                r.text = txt[1:]; r.font.color.rgb = Brand.ORANGE; r.font.bold = True
            else:
                r.text = txt; r.font.color.rgb = th.table_body_text
    return s


def _chart_in(slide, spec, left, top, width, height, th):
    """Pose un graphique natif charte dans un rectangle donne (reutilisable)."""
    ctype = {"bar": XL_CHART_TYPE.COLUMN_CLUSTERED, "line": XL_CHART_TYPE.LINE_MARKERS,
             "pie": XL_CHART_TYPE.PIE}.get(spec.get("chart_type", "bar"), XL_CHART_TYPE.COLUMN_CLUSTERED)
    data = CategoryChartData(); data.categories = spec["categories"]
    for serie in spec["series"]:
        data.add_series(serie["name"], serie["values"])
    gframe = slide.shapes.add_chart(ctype, Inches(left), Inches(top), Inches(width), Inches(height), data)
    chart = gframe.chart; chart.has_title = False
    palette = [Brand.ORANGE, th.text_muted, Brand.VERT, Brand.GRIS_CLAIR]
    for idx, ser in enumerate(chart.series):
        try:
            ser.format.fill.solid(); ser.format.fill.fore_color.rgb = palette[idx % len(palette)]
            ser.format.line.color.rgb = palette[idx % len(palette)]
        except Exception:
            pass
    if len(spec["series"]) > 1 and spec.get("chart_type") != "pie":
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.color.rgb = th.text; chart.legend.font.size = Pt(10); chart.legend.font.name = Brand.FONT
    for axis in (getattr(chart, "category_axis", None), getattr(chart, "value_axis", None)):
        if axis is None:
            continue
        try:
            axis.tick_labels.font.color.rgb = th.text
            axis.tick_labels.font.size = Pt(9); axis.tick_labels.font.name = Brand.FONT
            axis.format.line.color.rgb = th.text_muted
        except Exception:
            pass
    return chart


def _blocks_in(slide, blocks, left, top, width, height, th):
    """Blocs heading+body empiles dans un rectangle (colonne texte)."""
    n = max(1, len(blocks)); per = height / n
    y = top
    for b in blocks:
        _text(slide, b["heading"], left, y, width, 0.35, 14, Brand.ORANGE, bold=True)
        box = slide.shapes.add_textbox(Inches(left), Inches(y + 0.4), Inches(width), Inches(per - 0.45))
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]; p.line_spacing = 1.3
        _first_word_bold(p, b["body"], 12, th.text)
        y += per


def _cards_in(slide, cards, left, top, width, height, th):
    """Cartes (mini tuiles) empilees dans un rectangle (colonne cartes)."""
    n = max(1, len(cards)); gap = 0.2
    ch = (height - gap*(n-1)) / n; y = top
    for cd in cards:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(y), Inches(width), Inches(ch))
        _fill(card, th.tile); _round(card, 0.08)
        pad = 0.22
        yy = y + pad*0.7
        if cd.get("heading"):
            _text(slide, cd["heading"], left + pad, yy, width - 2*pad, 0.32, 12.5, Brand.ORANGE, bold=True)
            yy += 0.36
        if cd.get("body"):
            box = slide.shapes.add_textbox(Inches(left + pad), Inches(yy), Inches(width - 2*pad), Inches(y + ch - yy - pad*0.5))
            tf = box.text_frame; tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]; p.line_spacing = 1.2
            _first_word_bold(p, cd["body"], 10.5, th.text)
        y += ch + gap


_RATIO = {"half": 0.5, "third": 0.34, "two-thirds": 0.64}


def layout_chart(prs, c, th):
    s = _blank(prs); _solid_bg(s, th.bg); _title(s, c["title"], th)
    _chart_in(s, c, Brand.MARGIN, 2.0, Brand.SW - 2*Brand.MARGIN, 4.4, th)
    return s


def layout_split(prs, c, th):
    """Colonne visuelle (chart/image) + colonne texte (blocks ou cards)."""
    s = _blank(prs); _solid_bg(s, th.bg); _title(s, c["title"], th)
    top = 2.0; bottom = Brand.FOOTER_TOP - 0.2; height = bottom - top
    total_w = Brand.SW - 2*Brand.MARGIN; gutter = 0.45
    frac = _RATIO.get(c.get("ratio", "half"), 0.5)
    left = c.get("left", {}); right = c.get("right", {})
    swap = c.get("swap", False)  # visuel a droite si True
    vis_w = total_w * frac - gutter/2
    txt_w = total_w - vis_w - gutter
    if not swap:
        vis_x = Brand.MARGIN; txt_x = Brand.MARGIN + vis_w + gutter
    else:
        txt_x = Brand.MARGIN; vis_x = Brand.MARGIN + txt_w + gutter

    # Colonne visuelle
    if "chart" in left:
        _chart_in(s, left["chart"], vis_x, top, vis_w, height, th)
    elif left.get("image") and os.path.exists(left["image"]):
        s.shapes.add_picture(left["image"], Inches(vis_x), Inches(top), width=Inches(vis_w))
    else:
        ph = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(vis_x), Inches(top), Inches(vis_w), Inches(height))
        _fill(ph, th.tile)
        _text(s, "[ VISUEL ]", vis_x, top + height/2 - 0.2, vis_w, 0.4, 12, th.text_muted, align=PP_ALIGN.CENTER)

    # Colonne texte
    if "cards" in right:
        _cards_in(s, right["cards"], txt_x, top, txt_w, height, th)
    elif "blocks" in right:
        _blocks_in(s, right["blocks"], txt_x, top, txt_w, height, th)
    return s


def layout_capture(prs, c, th):
    s = _blank(prs); _solid_bg(s, th.bg); _title(s, c["title"], th)
    right_w = 6.6; right_left = Brand.SW - Brand.MARGIN - right_w
    img = c.get("image")
    if img and os.path.exists(img):
        s.shapes.add_picture(img, Inches(right_left), Inches(2.0), width=Inches(right_w))
    else:
        ph = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right_left), Inches(2.0), Inches(right_w), Inches(4.3))
        _fill(ph, th.tile)
        _text(s, "[ CAPTURE CLIENT ]", right_left, 4.0, right_w, 0.5, 13, th.text_muted, align=PP_ALIGN.CENTER)
    frictions = c.get("frictions", [])
    fx = Brand.MARGIN; fw = right_left - Brand.MARGIN - 0.4
    top0 = 2.0; bottom = Brand.FOOTER_TOP - 0.25
    n = max(1, len(frictions)); per = (bottom - top0) / n
    gap = min(0.18, per*0.22); fh = per - gap; y = top0
    for i, fr in enumerate(frictions):
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(fx), Inches(y), Inches(fw), Inches(fh))
        _fill(card, th.tile); _round(card, 0.12)
        _text(s, f"{i+1}. {fr['label']}", fx + 0.2, y, fw - 1.75, fh, 11, th.text, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        impact = fr.get("impact", "MOYEN").upper()
        col = {"ELEVE": Brand.ORANGE, "ÉLEVÉ": Brand.ORANGE, "MOYEN": th.badge_moyen, "FAIBLE": Brand.VERT}.get(impact, th.badge_moyen)
        bh = 0.34
        badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(fx + fw - 1.5), Inches(y + (fh-bh)/2), Inches(1.2), Inches(bh))
        _fill(badge, col); _round(badge, 0.3)
        tf = badge.text_frame; tf.word_wrap = False; tf.margin_top = tf.margin_bottom = 0; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        rr = pp.add_run(); rr.text = impact
        rr.font.size = Pt(8); rr.font.bold = True; rr.font.name = Brand.FONT; rr.font.color.rgb = Brand.BLANC
        y += per
    return s


def _tag(slide, left, top, width, label, value, th, accent=False):
    """Mini colonne label(caps gris) + valeur(bold) pour les recommandations."""
    _text(slide, label, left, top, width, 0.25, 8, th.text_muted, caps=True)
    _text(slide, value, left, top + 0.24, width, 0.3, 12, Brand.ORANGE if accent else th.text, bold=True)


def layout_recommendations(prs, c, th):
    """Rangees : icone + titre + impact/effort/priorite + actions (page 15)."""
    s = _blank(prs); _solid_bg(s, th.bg); _title(s, c["title"], th, c.get("subtitle"))
    items = c.get("items", [])
    top0 = 1.7 if not c.get("subtitle") else 2.05
    bottom = Brand.FOOTER_TOP - 0.2
    n = max(1, len(items)); gap = 0.22
    rowh = (bottom - top0 - gap*(n-1)) / n
    total_w = Brand.SW - 2*Brand.MARGIN
    y = top0
    for it in items:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(Brand.MARGIN), Inches(y),
                                  Inches(total_w), Inches(rowh))
        _fill(card, th.tile); _round(card, 0.06)
        pad = 0.28
        # icone
        isz = min(0.7, rowh * 0.5)
        _icon_square(s, Brand.MARGIN + pad, y + (rowh - isz)/2, isz, it.get("icon", "check-circle-2"),
                     th, orange=it.get("accent", False))
        # titre
        title_x = Brand.MARGIN + pad + isz + 0.25
        _text(s, it.get("title", ""), title_x, y, 3.4, rowh, 13.5, th.text, bold=True, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
        # tags impact / effort / priorite
        tags_x = title_x + 3.5
        tw = 1.15
        ty = y + (rowh - 0.54)/2
        cols = [("IMPACT", it.get("impact", "-"), False),
                ("EFFORT", it.get("effort", "-"), False),
                ("PRIORITE", str(it.get("priority", "-")), True)]
        for k, (lab, val, acc) in enumerate(cols):
            _tag(s, tags_x + k*tw, ty, tw, lab, val, th, accent=acc)
        # actions (puces)
        act_x = tags_x + 3*tw + 0.2
        act_w = Brand.SW - Brand.MARGIN - pad - act_x
        actions = it.get("actions", [])
        box = s.shapes.add_textbox(Inches(act_x), Inches(y + pad*0.5), Inches(act_w), Inches(rowh - pad))
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for j, a in enumerate(actions):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.line_spacing = 1.1
            r = p.add_run(); r.text = "- " + a
            r.font.size = Pt(10); r.font.name = Brand.FONT; r.font.color.rgb = th.text
        y += rowh + gap
    return s


def layout_audit(prs, c, th):
    """Audit CRO composite : resume + bandeau impact + frictions (gauche), capture (droite)."""
    s = _blank(prs); _solid_bg(s, th.bg); _title(s, c["title"], th, c.get("subtitle"))
    top = 2.05 if c.get("subtitle") else 1.8
    bottom = Brand.FOOTER_TOP - 0.2
    # colonne droite : capture
    right_w = 5.9; right_left = Brand.SW - Brand.MARGIN - right_w
    img = c.get("image")
    if img and os.path.exists(img):
        s.shapes.add_picture(img, Inches(right_left), Inches(top), width=Inches(right_w))
    else:
        ph = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right_left), Inches(top), Inches(right_w), Inches(bottom - top))
        _fill(ph, th.tile)
        _text(s, "[ CAPTURE CLIENT ]", right_left, (top + bottom)/2 - 0.2, right_w, 0.4, 12, th.text_muted, align=PP_ALIGN.CENTER)
    # colonne gauche
    lx = Brand.MARGIN; lw = right_left - Brand.MARGIN - 0.4
    y = top
    # resume executif
    if c.get("summary"):
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(lx), Inches(y), Inches(lw), Inches(1.15))
        _fill(card, th.tile); _round(card, 0.08)
        _icon_square(s, lx + 0.22, y + 0.2, 0.42, "trending-up", th)
        _text(s, "RESUME EXECUTIF", lx + 0.78, y + 0.2, lw - 1.0, 0.25, 9, th.text_muted, caps=True, bold=True)
        box = s.shapes.add_textbox(Inches(lx + 0.78), Inches(y + 0.48), Inches(lw - 1.0), Inches(0.6))
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]; p.line_spacing = 1.2
        r = p.add_run(); r.text = c["summary"]
        r.font.size = Pt(10.5); r.font.name = Brand.FONT; r.font.color.rgb = th.text
        y += 1.3
    # bandeau impact estime (mini KPIs)
    impacts = c.get("impact", [])
    if impacts:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(lx), Inches(y), Inches(lw), Inches(1.35))
        _fill(card, th.tile); _round(card, 0.08)
        _text(s, "IMPACT ESTIME", lx + 0.25, y + 0.16, lw - 0.5, 0.25, 9, th.text_muted, caps=True, bold=True)
        m = len(impacts); iw = (lw - 0.5) / m
        for k, imp in enumerate(impacts):
            ix = lx + 0.25 + k*iw
            _text(s, imp.get("label", ""), ix, y + 0.45, iw - 0.12, 0.35, 7.5, th.text_muted, caps=True, line_spacing=1.0)
            _text(s, imp.get("value", ""), ix, y + 0.74, iw - 0.12, 0.5, 11, Brand.ORANGE, bold=True, line_spacing=1.0)
        y += 1.5
    # frictions numerotees
    _text(s, "PRINCIPALES FRICTIONS IDENTIFIEES", lx, y, lw, 0.25, 9, th.text_muted, caps=True, bold=True)
    y += 0.32
    frictions = c.get("frictions", [])
    if frictions:
        n = len(frictions); avail = bottom - y
        per = avail / n; gap = min(0.12, per*0.18); fh = per - gap
        for i, fr in enumerate(frictions):
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(lx), Inches(y), Inches(lw), Inches(fh))
            _fill(card, th.tile); _round(card, 0.1)
            _text(s, f"{i+1}. {fr['label']}", lx + 0.18, y, lw - 1.5, fh, 10, th.text, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
            impact = fr.get("impact", "MOYEN").upper()
            col = {"ELEVE": Brand.ORANGE, "ÉLEVÉ": Brand.ORANGE, "MOYEN": th.badge_moyen, "FAIBLE": Brand.VERT}.get(impact, th.badge_moyen)
            bh = min(0.3, fh*0.7)
            badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(lx + lw - 1.25), Inches(y + (fh-bh)/2), Inches(1.05), Inches(bh))
            _fill(badge, col); _round(badge, 0.3)
            tf = badge.text_frame; tf.word_wrap = False; tf.margin_top = tf.margin_bottom = 0; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
            rr = pp.add_run(); rr.text = impact
            rr.font.size = Pt(7.5); rr.font.bold = True; rr.font.name = Brand.FONT; rr.font.color.rgb = Brand.BLANC
            y += per
    return s


def layout_closing(prs, c, th):
    s = _blank(prs); _gradient_bg(s, th.bg_grad_a, th.bg_grad_b)
    _logo(s, Brand.MARGIN, 2.35, 1.55, th, "principal")  # badge + wordmark THE OZ
    if c.get("headline"):
        _text(s, c["headline"], Brand.MARGIN, 4.15, 7, 0.4, 13, th.text_muted, bold=True)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(Brand.MARGIN), Inches(4.6), Inches(0.5), Pt(3))
    _fill(bar, Brand.ORANGE)
    _text(s, c.get("url", "www.the-oz.com"), Brand.MARGIN, 4.8, 6, 0.4, 14, Brand.LINK)
    return s


LAYOUTS = {
    "cover": layout_cover, "section": layout_section, "content": layout_content,
    "tiles": layout_tiles, "grid": layout_grid, "kpi": layout_kpi, "table": layout_table,
    "chart": layout_chart, "split": layout_split, "capture": layout_capture,
    "recommendations": layout_recommendations, "audit": layout_audit, "closing": layout_closing,
}
_NO_PAGE = {"cover", "closing"}


def build(schema: dict, output_path: str) -> str:
    """Point d'entree : schema dict -> fichier .pptx charte."""
    th = get_theme(schema.get("theme"))
    prs = Presentation()
    prs.slide_width = Inches(Brand.SW); prs.slide_height = Inches(Brand.SH)
    slides = schema.get("slides", [])
    total = len(slides)
    for idx, sl in enumerate(slides, start=1):
        layout = sl.get("layout")
        fn = LAYOUTS.get(layout)
        if fn is None:
            raise ValueError(f"Layout inconnu : {layout!r}")
        slide_theme = get_theme(sl.get("theme")) if sl.get("theme") else th
        s = fn(prs, sl, slide_theme)
        if layout in _NO_PAGE:
            _footer(s, slide_theme)
        else:
            _footer(s, slide_theme, page_no=idx, total=total)
    prs.save(output_path)
    return output_path
